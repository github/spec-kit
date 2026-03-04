"""
Spec-Kit vs Vibe Coding — PowerPoint Presentation Generator
Generates a comprehensive, visually rich PPTX presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ─── Brand Colors ────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x11, 0x17)   # GitHub dark
CARD_BG      = RGBColor(0x16, 0x1B, 0x22)   # Card surface
GREEN        = RGBColor(0x23, 0x86, 0x36)   # GitHub green
BLUE         = RGBColor(0x58, 0xA6, 0xFF)   # GitHub blue accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xC9, 0xD1, 0xD9)
MID_GRAY     = RGBColor(0x8B, 0x94, 0x9E)
ORANGE       = RGBColor(0xF7, 0x85, 0x16)
RED          = RGBColor(0xDA, 0x36, 0x33)
YELLOW       = RGBColor(0xF0, 0xD0, 0x44)
PURPLE       = RGBColor(0x79, 0x4C, 0xFF)
TEAL         = RGBColor(0x1A, 0x7F, 0x8A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]  # blank
    return prs.slides.add_slide(layout)


def bg(slide, color=DARK_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, color=CARD_BG, radius=None):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def rounded_box(slide, x, y, w, h, color=CARD_BG, line_color=None, line_width=None):
    """Add a rounded rectangle."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        5,  # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def txt_multi(slide, lines, x, y, w, h, size=16, bold=False,
              color=WHITE, align=PP_ALIGN.LEFT, line_spacing=1.2):
    """Add multi-line text with a text frame."""
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox


def accent_bar(slide, x, y, w=0.05, h=0.8, color=GREEN):
    """Thin vertical accent bar."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def circle(slide, x, y, d, color=GREEN):
    """Add a circle (ellipse)."""
    shape = slide.shapes.add_shape(
        9,  # OVAL
        Inches(x), Inches(y), Inches(d), Inches(d)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def arrow_right(slide, x, y, w=0.6, h=0.3, color=MID_GRAY):
    """Add a right-pointing arrow shape."""
    shape = slide.shapes.add_shape(
        13,  # RIGHT_ARROW
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def divider(slide, y, color=CARD_BG):
    """Horizontal rule."""
    shape = slide.shapes.add_shape(1, Inches(0), Inches(y), Inches(13.33), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def badge(slide, label, x, y, w=1.8, h=0.35, bg_color=GREEN, text_color=WHITE, size=11):
    b = rounded_box(slide, x, y, w, h, color=bg_color)
    txt(slide, label, x + 0.05, y + 0.02, w - 0.1, h - 0.04,
        size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)


def section_header(slide, title, subtitle=None, color=GREEN):
    """Top header band for section slides."""
    box(slide, 0, 0, 13.33, 1.1, color=color)
    txt(slide, title, 0.4, 0.12, 10, 0.65,
        size=32, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.72, 12, 0.35,
            size=14, color=RGBColor(0xE0, 0xE0, 0xE0))


def bullet_card(slide, title, bullets, x, y, w, h, title_color=BLUE,
                bg_c=CARD_BG, bullet_char="▶", size=13):
    rounded_box(slide, x, y, w, h, color=bg_c,
                line_color=RGBColor(0x30, 0x38, 0x42))
    txt(slide, title, x + 0.18, y + 0.15, w - 0.3, 0.38,
        size=16, bold=True, color=title_color)
    for i, b in enumerate(bullets):
        txt(slide, f"{bullet_char}  {b}", x + 0.18, y + 0.55 + i * 0.37,
            w - 0.3, 0.35, size=size, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title / Hero slide."""
    s = add_slide(prs)
    bg(s, DARK_BG)

    # Gradient-like top strip
    box(s, 0, 0, 13.33, 0.06, GREEN)
    box(s, 0, 0.06, 13.33, 0.06, RGBColor(0x1A, 0x72, 0x2A))

    # Large circle decoration
    circle(s, 9.5, -1.5, 6, color=RGBColor(0x12, 0x2A, 0x16))

    # Logo tag
    rounded_box(s, 0.5, 0.5, 1.8, 0.45, color=GREEN)
    txt(s, "SPEC-KIT", 0.5, 0.52, 1.8, 0.42,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Main heading
    txt(s, "Spec-Driven Development", 0.5, 1.3, 9, 0.85,
        size=46, bold=True, color=WHITE)
    txt(s, "vs  Vibe Coding", 0.5, 2.05, 7, 0.75,
        size=36, bold=False, color=BLUE)

    txt(s, "How Spec-Kit Transforms AI-Assisted Software Engineering",
        0.5, 2.95, 10, 0.5, size=18, color=LIGHT_GRAY)

    # Sub-line
    txt(s, "From Chaos to Clarity  ·  From Prompts to Executable Specifications",
        0.5, 3.55, 11, 0.4, size=14, italic=True, color=MID_GRAY)

    # Bottom stat pills
    for i, (val, label) in enumerate([
        ("17+", "AI Agents"), ("8", "Core Commands"),
        ("3", "Dev Phases"), ("0.1.13", "Latest Version")
    ]):
        bx = 0.5 + i * 3.1
        rounded_box(s, bx, 4.5, 2.7, 0.9, color=CARD_BG,
                    line_color=GREEN, line_width=1.5)
        txt(s, val, bx, 4.55, 2.7, 0.45,
            size=26, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        txt(s, label, bx, 5.0, 2.7, 0.36,
            size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Footer
    box(s, 0, 7.1, 13.33, 0.4, color=RGBColor(0x08, 0x0C, 0x10))
    txt(s, "github.com/github/spec-kit  ·  MIT License  ·  Open Source",
        0, 7.12, 13.33, 0.36, size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)


def slide_02_agenda(prs):
    s = add_slide(prs)
    bg(s, DARK_BG)
    section_header(s, "What We'll Cover", color=PURPLE)

    items = [
        ("01", "The Vibe Coding Crisis", "Why iterative, prompt-only dev breaks down at scale"),
        ("02", "Traditional SDLC History", "Where it came from and what it got right"),
        ("03", "Spec-Driven Development", "The methodology that bridges planning and AI execution"),
        ("04", "Spec-Kit Components", "8 commands, constitutional framework, multi-agent support"),
        ("05", "Workflow Comparison", "Step-by-step SDD vs Vibe Coding"),
        ("06", "Levels of Spec-Driven", "spec-as-source → spec-anchored → spec-first"),
        ("07", "Getting Started", "Install, init, and ship in minutes"),
    ]

    cols = [
        items[:4],
        items[4:],
    ]

    for ci, col in enumerate(cols):
        for ri, (num, title, desc) in enumerate(col):
            ox = 0.4 + ci * 6.5
            oy = 1.3 + ri * 1.45
            rounded_box(s, ox, oy, 6.1, 1.25,
                        color=CARD_BG, line_color=RGBColor(0x30, 0x38, 0x42))
            # Number badge
            circle(s, ox + 0.12, oy + 0.3, 0.55, color=GREEN)
            txt(s, num, ox + 0.12, oy + 0.33, 0.55, 0.48,
                size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            txt(s, title, ox + 0.82, oy + 0.12, 5.0, 0.38,
                size=15, bold=True, color=WHITE)
            txt(s, desc, ox + 0.82, oy + 0.52, 5.0, 0.65,
                size=11, color=MID_GRAY)


def slide_03_vibe_problems(prs):
    s = add_slide(prs)
    bg(s, DARK_BG)
    section_header(s, "The Vibe Coding Problem",
                   subtitle="Why 'just prompt the AI' falls apart at scale", color=RED)

    # Central diagram recreation
    # Vibe Coding loop visual
    box(s, 0.3, 1.25, 12.73, 5.9, color=RGBColor(0x12, 0x06, 0x06))
    box(s, 0.3, 1.25, 12.73, 0.04, color=RED)

    txt(s, '"Vibe Coding" Failure Loop', 0.7, 1.35, 6, 0.4,
        size=18, bold=True, color=RED)

    # Steps
    steps = [
        ("Write\nPrompt", BLUE, 1.2),
        ("AI Generates\nCode", ORANGE, 3.5),
        ("Test &\nEvaluate", YELLOW, 5.8),
        ("Not Right?\nEdit Prompt", RED, 8.1),
    ]

    for label, color, bx in steps:
        circle(s, bx, 2.1, 1.4, color=color)
        txt(s, label, bx, 2.1, 1.4, 1.4,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if bx < 8.1:
            arrow_right(s, bx + 1.45, 2.65, 0.6, 0.3,
                        color=MID_GRAY)

    # Loop-back arrow text
    txt(s, "↩  Iteration loop — can repeat dozens of times",
        1.0, 3.7, 8, 0.4, size=13, italic=True, color=ORANGE)

    # Problems list
    problems = [
        ("No Shared Memory", "AI forgets context between sessions — inconsistent decisions"),
        ("Spec Drift",        "Code diverges from original intent with every iteration"),
        ("Hidden Complexity", "Edge cases pile up invisibly until launch"),
        ("Time Waste",        "Hours of back-and-forth for tasks solvable in minutes with a spec"),
        ("No Contract",       "Nothing governs what 'done' means — quality is undefined"),
        ("Scale Failure",     "Works for demos; breaks for real products with multiple features"),
    ]

    for i, (title, desc) in enumerate(problems):
        col = i % 2
        row = i // 2
        ox = 0.5 + col * 6.4
        oy = 4.35 + row * 0.95
        rounded_box(s, ox, oy, 6.1, 0.82,
                    color=RGBColor(0x2A, 0x10, 0x10),
                    line_color=RED, line_width=0.8)
        txt(s, "✗", ox + 0.15, oy + 0.18, 0.35, 0.45,
            size=16, bold=True, color=RED)
        txt(s, title, ox + 0.55, oy + 0.08, 2.5, 0.32,
            size=13, bold=True, color=WHITE)
        txt(s, desc, ox + 0.55, oy + 0.4, 5.4, 0.38,
            size=11, color=MID_GRAY)


def slide_04_sdlc_history(prs):
    s = add_slide(prs)
    bg(s, DARK_BG)
    section_header(s, "Traditional SDLC — A Brief History",
                   subtitle="60+ years of software process evolution", color=TEAL)

    # Timeline bar
    box(s, 0.4, 1.35, 12.5, 0.06, color=TEAL)

    eras = [
        ("1960s–70s", "Waterfall", "Sequential phases,\nheavy docs", TEAL),
        ("1980s–90s", "Spiral / RAD", "Risk-driven,\nprototyping", BLUE),
        ("2000s",     "Agile / Scrum", "Iterative sprints,\nuser stories", GREEN),
        ("2010s",     "DevOps / CD", "CI/CD pipelines,\nInfrastructure-as-code", ORANGE),
        ("2020s+",    "AI-Assisted", "LLM co-pilots,\nSpec-Driven Dev", PURPLE),
    ]

    for i, (period, name, desc, color) in enumerate(eras):
        ox = 0.4 + i * 2.5
        # Dot on timeline
        circle(s, ox + 0.85, 1.15, 0.3, color=color)
        # Card
        rounded_box(s, ox, 1.65, 2.3, 2.1, color=CARD_BG,
                    line_color=color, line_width=1.5)
        txt(s, period, ox + 0.1, 1.72, 2.1, 0.3,
            size=10, color=MID_GRAY)
        txt(s, name, ox + 0.1, 2.0, 2.1, 0.42,
            size=14, bold=True, color=color)
        txt(s, desc, ox + 0.1, 2.45, 2.1, 0.75,
            size=11, color=LIGHT_GRAY)

    # SDLC phases circle illustration
    phases = ["Plan", "Analyze", "Design", "Implement", "Test", "Deploy", "Maintain"]
    cx, cy, r = 3.5, 5.5, 1.25
    import math
    for i, phase in enumerate(phases):
        angle = math.radians(-90 + i * (360 / len(phases)))
        px = cx + r * math.cos(angle) - 0.5
        py = cy + r * math.sin(angle) - 0.2
        c = [TEAL, BLUE, GREEN, ORANGE, RED, PURPLE, YELLOW][i]
        rounded_box(s, px, py, 1.0, 0.38, color=c)
        txt(s, phase, px, py + 0.02, 1.0, 0.34,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, "Classic SDLC Phases", 2.5, 6.9, 2.0, 0.35,
        size=11, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Similarities / Differences table
    box(s, 7.0, 3.95, 6.0, 3.3, color=RGBColor(0x10, 0x16, 0x1E))
    txt(s, "SDLC   ↔   Spec-Driven Dev", 7.1, 3.97, 5.8, 0.38,
        size=14, bold=True, color=WHITE)

    rows = [
        ("Similarities",  GREEN,  [
            "Phases: plan → design → build → test",
            "Requirements define the contract",
            "Quality gates before shipping",
        ]),
        ("Differences",   RED, [
            "SDD: specs are executable, not static docs",
            "SDD: AI is co-author throughout",
            "SDD: iteration is structured, not freeform",
        ]),
    ]

    ry = 4.42
    for header, color, items in rows:
        txt(s, header, 7.1, ry, 2.0, 0.3,
            size=12, bold=True, color=color)
        ry += 0.3
        for item in items:
            txt(s, f"• {item}", 7.1, ry, 5.7, 0.3,
                size=10, color=LIGHT_GRAY)
            ry += 0.3
        ry += 0.1


def slide_05_sdd_intro(prs):
    s = add_slide(prs)
    bg(s, DARK_BG)
    section_header(s, "Spec-Driven Development (SDD)",
                   subtitle="Specifications become the executable source of truth", color=GREEN)

    # Core definition card
    rounded_box(s, 0.4, 1.25, 12.5, 1.1, color=RGBColor(0x0C, 0x25, 0x14),
                line_color=GREEN, line_width=2)
    txt(s, "\"SDD is a methodology where specifications are not just documents — "
           "they are directly executable artifacts that drive code generation, "
           "quality assurance, and project governance.\"",
        0.65, 1.35, 12.0, 0.95,
        size=14, italic=True, color=WHITE)

    # 4 pillars
    pillars = [
        ("Specification\nas Source", "Specs generate code,\nnot just guide it", GREEN,    "📄"),
        ("Constitutional\nGovernance", "9 immutable articles\nshape every decision", BLUE,  "⚖"),
        ("Multi-Agent\nCompatibility", "Works with 17+ AI\ncoding assistants", ORANGE,    "🤖"),
        ("Quality\nBy Design", "Tests and checklists\nbaked into workflow", PURPLE,       "✔"),
    ]

    for i, (title, desc, color, icon) in enumerate(pillars):
        ox = 0.4 + i * 3.2
        rounded_box(s, ox, 2.55, 3.0, 2.0, color=CARD_BG,
                    line_color=color, line_width=1.5)
        # Icon circle
        circle(s, ox + 1.2, 2.65, 0.65, color=color)
        txt(s, icon, ox + 1.2, 2.65, 0.65, 0.65,
            size=22, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, title, ox + 0.1, 3.4, 2.8, 0.55,
            size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(s, desc, ox + 0.1, 3.95, 2.8, 0.55,
            size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # SDD vs Traditional
    txt(s, "SDD vs  Traditional Approaches", 0.4, 4.72, 6, 0.4,
        size=16, bold=True, color=WHITE)

    comparisons = [
        ("Requirement Docs",  "Static PDFs, rarely updated",   "Living specs that execute"),
        ("AI Role",           "Prompted ad-hoc",               "Constitutional co-author"),
        ("Quality Control",   "Manual review after coding",    "Checklist-driven, built-in"),
        ("Iteration",         "Unstructured guess-and-check",  "Governed by specification"),
    ]

    box(s, 0.4, 5.15, 12.5, 0.35, color=RGBColor(0x20, 0x27, 0x30))
    txt(s, "Aspect", 0.55, 5.18, 2.5, 0.3,
        size=11, bold=True, color=MID_GRAY)
    txt(s, "Traditional / Vibe Coding", 3.1, 5.18, 4.0, 0.3,
        size=11, bold=True, color=MID_GRAY)
    txt(s, "Spec-Driven Development", 7.2, 5.18, 5.0, 0.3,
        size=11, bold=True, color=MID_GRAY)

    for i, (aspect, trad, sdd) in enumerate(comparisons):
        oy = 5.55 + i * 0.42
        bg_c = CARD_BG if i % 2 == 0 else RGBColor(0x12, 0x18, 0x22)
        box(s, 0.4, oy, 12.5, 0.4, color=bg_c)
        txt(s, aspect, 0.55, oy + 0.05, 2.5, 0.32, size=11, bold=True, color=LIGHT_GRAY)
        txt(s, trad,   3.1,  oy + 0.05, 4.0, 0.32, size=11, color=RED)
        txt(s, sdd,    7.2,  oy + 0.05, 5.5, 0.32, size=11, color=GREEN)


def slide_06_speckit_overview(prs):
    s = add_slide(prs)
    bg(s, DARK_BG)
    section_header(s, "What is Spec-Kit?",
                   subtitle="An open-source toolkit for Spec-Driven Development", color=BLUE)

    # Left panel
    rounded_box(s, 0.3, 1.25, 5.8, 6.0, color=CARD_BG,
                line_color=BLUE, line_width=1.5)
    txt(s, "Core Capabilities", 0.5, 1.35, 5.4, 0.4,
        size=16, bold=True, color=BLUE)

    caps = [
        ("specify init", "Bootstrap any project with SDD in seconds"),
        ("8 Slash Commands", "From /speckit.specify to /speckit.implement"),
        ("17+ AI Agents", "Claude, Copilot, Gemini, Cursor, and more"),
        ("Extension System", "Modular add-ons without bloating core"),
        ("Constitutional FW", "9 articles that govern ALL decisions"),
        ("Cross-Artifact QA", "Analyze consistency across all specs"),
        ("Template Library", "Structured templates for specs, plans, tasks"),
    ]

    for i, (cap, desc) in enumerate(caps):
        oy = 1.85 + i * 0.72
        accent_bar(s, 0.42, oy, h=0.45, color=BLUE)
        txt(s, cap, 0.6, oy, 2.2, 0.28, size=13, bold=True, color=WHITE)
        txt(s, desc, 0.6, oy + 0.28, 5.3, 0.38, size=10, color=MID_GRAY)

    # Right panel — Architecture diagram
    txt(s, "Architecture", 6.5, 1.3, 6.5, 0.4,
        size=16, bold=True, color=WHITE)

    layers = [
        ("AI Agent Layer", "Claude · Copilot · Gemini · Cursor · 13 more", BLUE,   6.4, 1.8),
        ("Slash Commands", "/specify  /plan  /implement  /analyze  /tasks", GREEN,  6.4, 2.7),
        ("Template Engine", "spec.md · plan.md · tasks.md · constitution.md", ORANGE,6.4, 3.6),
        ("Extension System", "Community extensions · Version management", PURPLE,   6.4, 4.5),
        ("specify CLI", "specify init · specify upgrade · specify list", TEAL,     6.4, 5.4),
        ("Project Files", ".specify/  · memory/ · commands/", MID_GRAY,           6.4, 6.3),
    ]

    for (label, detail, color, lx, ly) in layers:
        rounded_box(s, lx, ly, 6.6, 0.72, color=RGBColor(0x0E, 0x15, 0x1E),
                    line_color=color, line_width=1.5)
        txt(s, label, lx + 0.18, ly + 0.06, 2.5, 0.3,
            size=13, bold=True, color=color)
        txt(s, detail, lx + 0.18, ly + 0.38, 6.2, 0.3,
            size=10, color=LIGHT_GRAY)
        if ly < 6.3:
            # Connector arrow down
            shape = s.shapes.add_shape(
                1, Inches(lx + 3.1), Inches(ly + 0.73),
                Inches(0.04), Inches(0.2)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = MID_GRAY
            shape.line.fill.background()


def slide_07_eight_commands(prs):
    s = add_slide(prs)
    bg(s, DARK_BG)
    section_header(s, "The 8 Core Spec-Kit Commands",
                   subtitle="A complete lifecycle from vision to verified implementation", color=GREEN)

    commands = [
        ("/speckit.constitution", "Establish governing principles",
         "Define the project's immutable rules: coding standards, architecture decisions, library choices.",
         GREEN, "⚖"),
        ("/speckit.specify", "Define feature requirements",
         "Create structured specs with user stories, acceptance criteria, and constraints.",
         BLUE, "📋"),
        ("/speckit.clarify", "Resolve ambiguities",
         "Structured Q&A workflow to eliminate assumptions before any code is written.",
         ORANGE, "❓"),
        ("/speckit.plan", "Generate implementation plans",
         "Turn requirements into concrete technical design documents with specific to-dos.",
         YELLOW, "🗺"),
        ("/speckit.analyze", "Cross-artifact consistency",
         "Validate alignment between specs, plans, and code across the entire project.",
         PURPLE, "🔍"),
        ("/speckit.tasks", "Break plans into tasks",
         "Convert design documents into actionable, sequenced implementation tasks.",
         TEAL, "✅"),
        ("/speckit.implement", "Execute all tasks",
         "LLM agent writes code guided by design doc — less ambiguity, better output.",
         RED, "⚡"),
        ("/speckit.checklist", "Generate quality checklists",
         "Custom QA checklists derived from your requirements for validation and review.",
         RGBColor(0xFF, 0x79, 0xC6), "☑"),
    ]

    for i, (cmd, subtitle, desc, color, icon) in enumerate(commands):
        col = i % 2
        row = i // 2
        ox = 0.35 + col * 6.5
        oy = 1.3 + row * 1.48

        rounded_box(s, ox, oy, 6.2, 1.35, color=CARD_BG,
                    line_color=color, line_width=1.5)
        # Icon circle
        circle(s, ox + 0.18, oy + 0.38, 0.6, color=color)
        txt(s, icon, ox + 0.18, oy + 0.38, 0.6, 0.6,
            size=18, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, cmd, ox + 0.95, oy + 0.08, 5.0, 0.35,
            size=13, bold=True, color=color)
        txt(s, subtitle, ox + 0.95, oy + 0.43, 5.1, 0.28,
            size=11, bold=True, color=WHITE)
        txt(s, desc, ox + 0.95, oy + 0.72, 5.1, 0.55,
            size=10, color=MID_GRAY)


def slide_08_constitution(prs):
    s = add_slide(prs)
    bg(s, DARK_BG)
    section_header(s, "The Constitutional Framework",
                   subtitle="9 immutable articles that govern every development decision", color=ORANGE)

    # Constitution document visual
    rounded_box(s, 0.3, 1.25, 4.2, 5.9, color=RGBColor(0x18, 0x12, 0x05),
                line_color=ORANGE, line_width=2)
    txt(s, "📜  Project Constitution", 0.5, 1.35, 3.8, 0.42,
        size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(s, ".specify/memory/\nconstitution.md", 0.5, 1.82, 3.8, 0.55,
        size=11, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    articles_left = [
        ("I",   "Library-First", "All features start as libraries"),
        ("II",  "CLI Mandate",   "Text I/O, JSON support required"),
        ("III", "Test-First",    "Tests written before implementation"),
        ("IV",  "Docs Living",   "Specs updated with every change"),
        ("V",   "Single Source", "One source of truth per artifact"),
    ]

    for i, (num, title, desc) in enumerate(articles_left):
        oy = 2.5 + i * 0.88
        txt(s, f"Art. {num}", 0.45, oy, 0.85, 0.28,
            size=10, bold=True, color=ORANGE)
        txt(s, title, 1.35, oy, 2.1, 0.28,
            size=11, bold=True, color=WHITE)
        txt(s, desc, 0.45, oy + 0.28, 3.9, 0.52,
            size=10, color=MID_GRAY)

    # Right side — articles VI-IX + explanation
    articles_right = [
        ("VI",   "Semantic Versioning", "All specs carry version numbers"),
        ("VII",  "Simplicity",           "Max 3 projects initial implementation"),
        ("VIII", "Anti-Abstraction",     "Use frameworks directly, no wrappers"),
        ("IX",   "Integration Testing",  "Real databases, not mocks"),
    ]

    txt(s, "Remaining Articles", 4.8, 1.3, 5, 0.4,
        size=14, bold=True, color=WHITE)

    for i, (num, title, desc) in enumerate(articles_right):
        oy = 1.8 + i * 0.82
        rounded_box(s, 4.8, oy, 4.5, 0.7, color=RGBColor(0x18, 0x12, 0x05),
                    line_color=RGBColor(0x50, 0x30, 0x10), line_width=0.8)
        txt(s, f"Article {num}", 4.95, oy + 0.06, 1.0, 0.25,
            size=10, bold=True, color=ORANGE)
        txt(s, title, 6.0, oy + 0.06, 3.1, 0.28,
            size=12, bold=True, color=WHITE)
        txt(s, desc, 4.95, oy + 0.38, 4.2, 0.28,
            size=10, color=MID_GRAY)

    # Why it matters
    rounded_box(s, 4.8, 5.12, 8.2, 2.0, color=RGBColor(0x0C, 0x16, 0x0C),
                line_color=GREEN, line_width=1.5)
    txt(s, "Why the Constitution Matters", 5.0, 5.2, 7, 0.38,
        size=15, bold=True, color=GREEN)
    reasons = [
        "▶  Eliminates contradictory decisions across a long project lifetime",
        "▶  Gives AI agents a stable, trustworthy contract to code against",
        "▶  Reduces 'creative drift' where each session forgets prior choices",
        "▶  Creates onboarding documentation automatically for new contributors",
    ]
    for i, r in enumerate(reasons):
        txt(s, r, 5.0, 5.65 + i * 0.34, 7.8, 0.3,
            size=11, color=LIGHT_GRAY)


def slide_09_levels(prs):
    s = add_slide(prs)
    bg(s, DARK_BG)
    section_header(s, "Levels of Spec-Driven Development",
                   subtitle="Choose your commitment level — all are valid starting points", color=TEAL)

    levels = [
        (
            "spec-as-source",
            "Entry Level",
            ORANGE,
            [
                "Spec created alongside code",
                "AI can write code without spec guidance",
                "Spec updated after code changes",
                "Best for: exploratory prototypes",
            ],
            "Creation: Spec + Code (parallel)\nEvolution: Code first, spec follows",
        ),
        (
            "spec-anchored",
            "Intermediate",
            BLUE,
            [
                "Spec guides initial implementation",
                "Spec updated before major changes",
                "AI always references spec for context",
                "Best for: production features",
            ],
            "Creation: Spec → Code\nEvolution: Edit spec → Update code",
        ),
        (
            "spec-first",
            "Advanced",
            GREEN,
            [
                "Spec is the single source of truth",
                "No code changes without spec change",
                "Old spec discarded, new spec written",
                "Best for: mission-critical systems",
            ],
            "Creation: Spec → Code\nEvolution: New spec → New code",
        ),
    ]

    for i, (name, level_label, color, bullets, flow) in enumerate(levels):
        ox = 0.4 + i * 4.3
        # Main card
        rounded_box(s, ox, 1.3, 4.1, 5.8, color=CARD_BG,
                    line_color=color, line_width=2)

        # Level badge
        badge(s, level_label, ox + 0.1, 1.38, w=1.5, h=0.34,
              bg_color=color, size=10)

        txt(s, name, ox + 0.15, 1.82, 3.8, 0.45,
            size=17, bold=True, color=color)

        # Flow description box
        rounded_box(s, ox + 0.1, 2.32, 3.9, 0.65,
                    color=RGBColor(0x10, 0x18, 0x25),
                    line_color=RGBColor(0x25, 0x35, 0x45))
        txt(s, flow, ox + 0.2, 2.37, 3.7, 0.58,
            size=10, italic=True, color=MID_GRAY)

        # Bullets
        for j, b in enumerate(bullets):
            txt(s, f"✓  {b}", ox + 0.15, 3.1 + j * 0.58, 3.8, 0.5,
                size=12, color=LIGHT_GRAY)

        # Vertical position indicator
        circle(s, ox + 1.85, 6.85, 0.25, color=color)

    # Arrow showing progression
    txt(s, "Increasing rigor & spec fidelity  →",
        0.4, 7.05, 8, 0.35,
        size=13, italic=True, color=MID_GRAY)


def slide_10_workflow_compare(prs):
    s = add_slide(prs)
    bg(s, DARK_BG)
    section_header(s, "Workflow Side-by-Side",
                   subtitle="Vibe Coding vs Spec-Driven Development", color=PURPLE)

    # Column headers
    box(s, 0.3, 1.25, 6.1, 0.55, color=RGBColor(0x25, 0x10, 0x10))
    txt(s, "⚡  Vibe Coding", 0.5, 1.3, 5.8, 0.45,
        size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)

    box(s, 6.9, 1.25, 6.1, 0.55, color=RGBColor(0x0C, 0x22, 0x12))
    txt(s, "✓  Spec-Driven Dev", 6.9, 1.3, 6.1, 0.45,
        size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # VS divider
    circle(s, 6.2, 1.32, 0.58, color=PURPLE)
    txt(s, "VS", 6.2, 1.32, 0.58, 0.58,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    vibe_steps = [
        ("1. Write a prompt", "Hope the AI understands intent"),
        ("2. Get boilerplate", "AI generates generic code"),
        ("3. Test manually", "Discover edge cases the hard way"),
        ("4. Prompt again", "Describe what's wrong, repeat"),
        ("5. Accumulate debt", "Each iteration adds hidden complexity"),
        ("6. Ship & pray", "No contract for 'done'"),
    ]

    sdd_steps = [
        ("1. /speckit.constitution", "Establish immutable project principles"),
        ("2. /speckit.specify", "Write structured spec with criteria"),
        ("3. /speckit.clarify", "Resolve every ambiguity upfront"),
        ("4. /speckit.plan", "Generate approved design document"),
        ("5. /speckit.tasks + /speckit.implement", "Execute against design"),
        ("6. /speckit.checklist", "Verify against requirements — ship"),
    ]

    for i, ((vt, vd), (st, sd)) in enumerate(zip(vibe_steps, sdd_steps)):
        oy = 1.95 + i * 0.84

        # Vibe step
        rounded_box(s, 0.3, oy, 6.1, 0.75, color=RGBColor(0x1E, 0x10, 0x10),
                    line_color=RGBColor(0x5A, 0x1A, 0x1A), line_width=0.8)
        txt(s, vt, 0.5, oy + 0.06, 3.5, 0.28,
            size=12, bold=True, color=RED)
        txt(s, vd, 0.5, oy + 0.38, 5.8, 0.32,
            size=10, color=MID_GRAY)

        # SDD step
        rounded_box(s, 6.9, oy, 6.1, 0.75, color=RGBColor(0x0A, 0x1E, 0x12),
                    line_color=RGBColor(0x1A, 0x5A, 0x28), line_width=0.8)
        txt(s, st, 7.1, oy + 0.06, 5.5, 0.28,
            size=12, bold=True, color=GREEN)
        txt(s, sd, 7.1, oy + 0.38, 5.8, 0.32,
            size=10, color=MID_GRAY)


def slide_11_agents(prs):
    s = add_slide(prs)
    bg(s, DARK_BG)
    section_header(s, "Multi-Agent Support",
                   subtitle="17+ AI coding assistants — one methodology", color=BLUE)

    agents = [
        ("Claude Code",    "Anthropic",   BLUE),
        ("GitHub Copilot", "GitHub",      GREEN),
        ("Gemini CLI",     "Google",      YELLOW),
        ("Cursor",         "Cursor",      PURPLE),
        ("Kiro",           "Kiro",        TEAL),
        ("Windsurf",       "Codeium",     ORANGE),
        ("Codex CLI",      "OpenAI",      RGBColor(0x74, 0xAA, 0x9C)),
        ("Roo Code",       "Roo",         RED),
        ("Qwen Code",      "Alibaba",     RGBColor(0xFF, 0x6A, 0x00)),
        ("IBM Bob",        "IBM",         BLUE),
        ("opencode",       "Apache",      GREEN),
        ("SHAI",           "OVHcloud",    PURPLE),
        ("Auggie",         "Auggie",      TEAL),
        ("Agy",            "Antigravity", ORANGE),
        ("CodeBuddy",      "CodeBuddy",   RGBColor(0xFF, 0x82, 0xB4)),
        ("Qoder",          "Qoder",       YELLOW),
        ("Generic Agent",  "Bring Your Own", MID_GRAY),
    ]

    cols = 5
    for i, (name, company, color) in enumerate(agents):
        col = i % cols
        row = i // cols
        ox = 0.4 + col * 2.55
        oy = 1.35 + row * 1.35
        rounded_box(s, ox, oy, 2.3, 1.1, color=CARD_BG,
                    line_color=color, line_width=1.5)
        circle(s, ox + 0.82, oy + 0.1, 0.42, color=color)
        txt(s, name[:2].upper(), ox + 0.82, oy + 0.1, 0.42, 0.42,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, name, ox + 0.05, oy + 0.6, 2.2, 0.28,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, company, ox + 0.05, oy + 0.88, 2.2, 0.2,
            size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Usage command at bottom
    rounded_box(s, 0.4, 6.5, 12.5, 0.7, color=RGBColor(0x0A, 0x12, 0x20),
                line_color=BLUE, line_width=1.5)
    txt(s, "specify init my-project --ai claude",
        1.0, 6.6, 8, 0.5,
        size=16, bold=True, color=GREEN)
    txt(s, "# or: --ai copilot  --ai gemini  --ai cursor  --ai kiro  ...",
        1.0, 6.85, 10, 0.3, size=10, italic=True, color=MID_GRAY)


def slide_12_getting_started(prs):
    s = add_slide(prs)
    bg(s, DARK_BG)
    section_header(s, "Getting Started with Spec-Kit",
                   subtitle="From zero to spec-driven in 4 steps", color=GREEN)

    steps = [
        ("01  Install", GREEN,
         "uv tool install specify-cli \\\n  --from git+https://github.com/github/spec-kit.git",
         "Requires Python 3.11+. Works on macOS, Windows, Linux."),
        ("02  Initialize", BLUE,
         "specify init my-awesome-project --ai claude\n# Or: specify init . --ai copilot",
         "Creates .specify/ folder with templates, commands, and agent config."),
        ("03  Spec & Plan", ORANGE,
         "/speckit.specify Build a photo album app with drag-and-drop\n/speckit.plan Use Vite, vanilla JS, SQLite",
         "AI generates structured spec and implementation plan — reviewed before any code."),
        ("04  Implement", PURPLE,
         "/speckit.tasks   # → actionable task list\n/speckit.implement  # → AI writes code against design",
         "Tests generated from requirements. /speckit.checklist verifies quality before ship."),
    ]

    for i, (title, color, code, note) in enumerate(steps):
        oy = 1.3 + i * 1.5
        rounded_box(s, 0.3, oy, 12.7, 1.35, color=CARD_BG,
                    line_color=color, line_width=2)

        # Step number badge
        rounded_box(s, 0.35, oy + 0.07, 1.55, 0.5, color=color)
        txt(s, title[:2], 0.35, oy + 0.07, 1.55, 0.5,
            size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        txt(s, title[4:], 2.05, oy + 0.07, 3, 0.38,
            size=15, bold=True, color=color)

        # Code box
        box(s, 5.0, oy + 0.08, 7.7, 0.75, color=RGBColor(0x08, 0x0D, 0x14))
        txt(s, code, 5.1, oy + 0.11, 7.5, 0.7,
            size=10, bold=True, color=GREEN)

        txt(s, note, 2.05, oy + 0.52, 2.75, 0.75,
            size=10, italic=True, color=MID_GRAY)


def slide_13_results(prs):
    s = add_slide(prs)
    bg(s, DARK_BG)
    section_header(s, "Real-World Benefits",
                   subtitle="What teams experience when they switch to Spec-Driven Development", color=PURPLE)

    # Metric cards
    metrics = [
        ("Fewer\nRework Cycles", "Up to 80%", "Specs resolve ambiguity before code is written", GREEN),
        ("Faster\nOnboarding", "5×", "New contributors understand the project via constitution + specs", BLUE),
        ("Bug Discovery", "Earlier", "Requirements-driven tests catch issues before integration", ORANGE),
        ("AI Consistency", "Higher", "Constitutional governance keeps AI decisions aligned across sessions", PURPLE),
    ]

    for i, (label, val, desc, color) in enumerate(metrics):
        ox = 0.4 + i * 3.2
        rounded_box(s, ox, 1.3, 3.0, 2.5, color=CARD_BG,
                    line_color=color, line_width=2)
        txt(s, val, ox, 1.4, 3.0, 0.85,
            size=38, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(s, label, ox, 2.25, 3.0, 0.6,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc, ox + 0.1, 2.88, 2.8, 0.85,
            size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Use-case cards
    txt(s, "Ideal Use Cases", 0.4, 4.0, 5, 0.4,
        size=16, bold=True, color=WHITE)

    cases = [
        ("Greenfield Projects",    "Start from zero with AI-generated architecture, specs, and tests",      GREEN),
        ("Legacy Modernization",   "Spec existing code first, then refactor guided by those specs",          BLUE),
        ("Multi-Team Features",    "Shared constitution aligns distributed teams building related features",  ORANGE),
        ("Rapid Prototyping",      "Spec-as-source level lets you move fast without losing context",         PURPLE),
        ("Regulated Industries",   "Specs become compliance artifacts; checklist ensures audit readiness",   TEAL),
        ("Open Source Projects",   "Constitution + specs onboard contributors without synchronous meetings", RED),
    ]

    for i, (title, desc, color) in enumerate(cases):
        col = i % 2
        row = i // 2
        ox = 0.4 + col * 6.5
        oy = 4.5 + row * 0.92
        rounded_box(s, ox, oy, 6.2, 0.78, color=CARD_BG,
                    line_color=color, line_width=1.2)
        accent_bar(s, ox + 0.08, oy + 0.12, h=0.55, color=color)
        txt(s, title, ox + 0.28, oy + 0.08, 2.5, 0.3,
            size=12, bold=True, color=WHITE)
        txt(s, desc, ox + 0.28, oy + 0.4, 5.7, 0.32,
            size=10, color=MID_GRAY)


def slide_14_cta(prs):
    s = add_slide(prs)
    bg(s, DARK_BG)

    # Gradient strip
    box(s, 0, 0, 13.33, 0.08, color=GREEN)

    # Large decorative circle
    circle(s, 8.5, -2.0, 8.0, color=RGBColor(0x0C, 0x22, 0x12))

    txt(s, "Start Building Smarter Today", 0.6, 0.5, 9, 0.8,
        size=42, bold=True, color=WHITE)
    txt(s, "Stop vibe coding. Start spec-driving.",
        0.6, 1.35, 8, 0.5, size=24, color=GREEN)
    txt(s, "Spec-Kit is free, open-source, and works with every major AI coding agent.",
        0.6, 2.0, 9.5, 0.4, size=15, color=LIGHT_GRAY)

    # CTA boxes
    ctas = [
        ("⭐  Star on GitHub", "github.com/github/spec-kit", GREEN,    0.6),
        ("📦  Install Now",    "specify init . --ai <agent>", BLUE,    4.0),
        ("📖  Read the Docs",  "Spec-Driven methodology guide", ORANGE, 7.4),
    ]

    for label, sub, color, ox in ctas:
        rounded_box(s, ox, 2.7, 3.0, 1.2, color=color)
        txt(s, label, ox + 0.1, 2.82, 2.8, 0.5,
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, sub, ox + 0.1, 3.3, 2.8, 0.5,
            size=11, color=RGBColor(0xE0, 0xFF, 0xE8), align=PP_ALIGN.CENTER)

    # Quote
    rounded_box(s, 0.5, 4.15, 12.3, 1.05, color=CARD_BG,
                line_color=GREEN, line_width=1.5)
    txt(s, '"Specifications are not overhead — they are the highest-leverage investment\n'
           ' you can make in AI-assisted software development."',
        0.75, 4.25, 11.8, 0.85,
        size=15, italic=True, color=WHITE)

    # Bottom links
    box(s, 0, 5.45, 13.33, 2.05, color=RGBColor(0x08, 0x0C, 0x10))
    links = [
        ("🌐 GitHub",            "github.com/github/spec-kit"),
        ("🐍 Install (pip/uv)",  "specify-cli on PyPI"),
        ("📚 Docs",              "spec-kit documentation"),
        ("💬 Community",         "GitHub Discussions"),
        ("📜 Methodology",       "spec-driven.md in repo"),
    ]
    for i, (icon, url) in enumerate(links):
        ox = 0.5 + i * 2.58
        txt(s, icon, ox, 5.55, 2.3, 0.35,
            size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        txt(s, url, ox, 5.9, 2.3, 0.38,
            size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    txt(s, "MIT License  ·  Open Source  ·  Contributions Welcome",
        0, 6.85, 13.33, 0.4,
        size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    slide_01_title(prs)
    print("  [ok] Slide 01 - Title")
    slide_02_agenda(prs)
    print("  [ok] Slide 02 - Agenda")
    slide_03_vibe_problems(prs)
    print("  [ok] Slide 03 - Vibe Coding Problems")
    slide_04_sdlc_history(prs)
    print("  [ok] Slide 04 - SDLC History")
    slide_05_sdd_intro(prs)
    print("  [ok] Slide 05 - SDD Introduction")
    slide_06_speckit_overview(prs)
    print("  [ok] Slide 06 - Spec-Kit Overview")
    slide_07_eight_commands(prs)
    print("  [ok] Slide 07 - 8 Core Commands")
    slide_08_constitution(prs)
    print("  [ok] Slide 08 - Constitutional Framework")
    slide_09_levels(prs)
    print("  [ok] Slide 09 - Levels of Spec-Driven")
    slide_10_workflow_compare(prs)
    print("  [ok] Slide 10 - Workflow Comparison")
    slide_11_agents(prs)
    print("  [ok] Slide 11 - Multi-Agent Support")
    slide_12_getting_started(prs)
    print("  [ok] Slide 12 - Getting Started")
    slide_13_results(prs)
    print("  [ok] Slide 13 - Results & Use Cases")
    slide_14_cta(prs)
    print("  [ok] Slide 14 - Call to Action")

    out = "presentations/spec-kit-vs-vibe-coding.pptx"
    prs.save(out)
    print(f"\nPresentation saved -> {out}")
    print(f"   {len(prs.slides)} slides  |  Widescreen 13.33\" x 7.5\"")


if __name__ == "__main__":
    main()
